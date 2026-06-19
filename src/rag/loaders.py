import importlib.util
import io
from collections.abc import Callable
from pathlib import Path

from src.core.exceptions import DocumentParseError, RagValidationError
from src.rag.models import (
    CONTENT_KIND_BY_FILE_TYPE,
    DocumentFileType,
    LoadedDocument,
)


SUPPORTED_EXTENSIONS = tuple(sorted(file_type.value for file_type in DocumentFileType))

_DOCLING_AVAILABLE = importlib.util.find_spec("docling") is not None


def load_document(*, source_name: str | None, raw_content: bytes) -> LoadedDocument:
    """
    Converts uploaded bytes into the document used by ingest.

    The caller reads the file bytes; this function validates them, extracts plain
    text with the parser its file type requires, and tags the result with the
    chunking strategy (``content_kind``) that text needs. Text-like types are
    decoded as UTF-8; binary types (PDF) are parsed into text.

    Binary formats (PDF, DOCX, PPTX, XLSX) use Docling when installed for
    ML-powered layout and table extraction. When Docling is absent the parser
    falls back to lightweight libraries (pypdf, python-docx, python-pptx,
    openpyxl) that extract plain text without structural analysis.
    """
    normalized_source_name = _normalize_source_name(source_name)
    file_type = _resolve_file_type(normalized_source_name)
    content = _PARSERS[file_type](normalized_source_name, raw_content)

    if not content.strip():
        raise RagValidationError(
            "document content must not be empty",
            details={"source_name": normalized_source_name},
        )

    return LoadedDocument(
        source_name=normalized_source_name,
        content=content,
        content_kind=CONTENT_KIND_BY_FILE_TYPE[file_type],
    )


def _resolve_file_type(source_name: str) -> DocumentFileType:
    extension = Path(source_name).suffix.lower()
    try:
        return DocumentFileType(extension)
    except ValueError as exc:
        raise RagValidationError(
            "Unsupported document type",
            details={
                "source_name": source_name,
                "supported_extensions": list(SUPPORTED_EXTENSIONS),
            },
        ) from exc


def _parse_utf8_text(source_name: str, raw_content: bytes) -> str:
    try:
        return raw_content.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise RagValidationError(
            "Document must be UTF-8 encoded text",
            details={"source_name": source_name},
        ) from exc


# ---------------------------------------------------------------------------
# Lightweight fallback parsers (used when docling is not installed)
# ---------------------------------------------------------------------------


def _parse_pdf_lightweight(source_name: str, raw_content: bytes) -> str:
    from pypdf import PdfReader

    try:
        reader = PdfReader(io.BytesIO(raw_content))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as exc:
        raise DocumentParseError(
            details={"source_name": source_name, "error": str(exc)}
        ) from exc


def _parse_docx_lightweight(source_name: str, raw_content: bytes) -> str:
    from docx import Document

    try:
        doc = Document(io.BytesIO(raw_content))
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception as exc:
        raise DocumentParseError(
            details={"source_name": source_name, "error": str(exc)}
        ) from exc


def _parse_pptx_lightweight(source_name: str, raw_content: bytes) -> str:
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(raw_content))
        texts = [
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        ]
        return "\n".join(texts)
    except Exception as exc:
        raise DocumentParseError(
            details={"source_name": source_name, "error": str(exc)}
        ) from exc


def _parse_xlsx_lightweight(source_name: str, raw_content: bytes) -> str:
    from openpyxl import load_workbook

    try:
        wb = load_workbook(io.BytesIO(raw_content), read_only=True, data_only=True)
        rows = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                rows.append(",".join("" if c is None else str(c) for c in row))
        return "\n".join(rows)
    except Exception as exc:
        raise DocumentParseError(
            details={"source_name": source_name, "error": str(exc)}
        ) from exc


# ---------------------------------------------------------------------------
# Docling parser (ML-powered; only instantiated when docling is installed)
# ---------------------------------------------------------------------------


class DoclingParser:
    """Lazy-loading Docling converter — initialised once per process, reused per request.

    ``to_markdown`` is used for prose formats (PDF, DOCX, PPTX).
    ``to_csv_tables`` is used for spreadsheets (XLSX) so the output feeds
    directly into ``TableChunker`` without row-boundary splits.
    """

    def __init__(self) -> None:
        self._converter = None

    @property
    def converter(self):
        if self._converter is None:
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import PdfPipelineOptions
            from docling.document_converter import DocumentConverter, PdfFormatOption

            pipeline_opts = PdfPipelineOptions()

            self._converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_opts)
                }
            )
        return self._converter

    def to_markdown(self, source_name: str, raw_content: bytes) -> str:
        return self._convert(source_name, raw_content).document.export_to_markdown()

    def to_csv_tables(self, source_name: str, raw_content: bytes) -> str:
        result = self._convert(source_name, raw_content)
        tables = [
            table.export_to_dataframe().to_csv(index=False)
            for table in result.document.tables
        ]
        return "\n".join(tables)

    def _convert(self, source_name: str, raw_content: bytes):
        from docling.datamodel.base_models import ConversionStatus
        from docling.datamodel.document import DocumentStream

        try:
            stream = DocumentStream(name=source_name, stream=io.BytesIO(raw_content))
            result = self.converter.convert(stream)
        except Exception as exc:
            raise DocumentParseError(
                details={"source_name": source_name, "error": str(exc)}
            ) from exc

        if result.status == ConversionStatus.FAILURE:
            raise DocumentParseError(
                details={"source_name": source_name, "status": result.status.value}
            )

        return result


# ---------------------------------------------------------------------------
# Parser selection — Docling when available, lightweight fallback otherwise
# ---------------------------------------------------------------------------

if _DOCLING_AVAILABLE:
    _docling_parser = DoclingParser()
    _binary_parsers: dict[DocumentFileType, Callable[[str, bytes], str]] = {
        DocumentFileType.PDF: _docling_parser.to_markdown,
        DocumentFileType.DOCX: _docling_parser.to_markdown,
        DocumentFileType.PPTX: _docling_parser.to_markdown,
        DocumentFileType.XLSX: _docling_parser.to_csv_tables,
    }
else:
    _binary_parsers = {
        DocumentFileType.PDF: _parse_pdf_lightweight,
        DocumentFileType.DOCX: _parse_docx_lightweight,
        DocumentFileType.PPTX: _parse_pptx_lightweight,
        DocumentFileType.XLSX: _parse_xlsx_lightweight,
    }

_PARSERS: dict[DocumentFileType, Callable[[str, bytes], str]] = {
    DocumentFileType.MARKDOWN: _parse_utf8_text,
    DocumentFileType.TEXT: _parse_utf8_text,
    DocumentFileType.CSV: _parse_utf8_text,
    **_binary_parsers,
}


def _normalize_source_name(source_name: str | None) -> str:
    if source_name is None or not source_name.strip():
        raise RagValidationError("source_name must not be empty")

    stripped = source_name.strip()
    normalized_path = stripped.replace("\\", "/")
    filename = normalized_path.rsplit("/", maxsplit=1)[-1]
    if not filename:
        raise RagValidationError("source_name must not be empty")

    return filename
