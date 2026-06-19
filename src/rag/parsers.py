import importlib.util
import io
from pathlib import Path

from src.core.exceptions import DocumentParseError


_DOCLING_AVAILABLE = importlib.util.find_spec("docling") is not None


class DoclingParser:
    """ML-powered parser using Docling (layout analysis, table detection).

    Lazy-initialised once per process and reused across requests.
    """

    def __init__(self) -> None:
        self._converter = None

    @property
    def converter(self):
        if self._converter is None:
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import PdfPipelineOptions
            from docling.document_converter import DocumentConverter, PdfFormatOption

            self._converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(
                        pipeline_options=PdfPipelineOptions()
                    )
                }
            )
        return self._converter

    def to_markdown(self, source_name: str, raw_content: bytes) -> str:
        return self._convert(source_name, raw_content).document.export_to_markdown()

    def to_csv_tables(self, source_name: str, raw_content: bytes) -> str:
        result = self._convert(source_name, raw_content)
        tables = [
            table.export_to_dataframe().to_csv(index=False)
            for table in result.document.tables
        ]
        return "\n".join(tables)

    def _convert(self, source_name: str, raw_content: bytes):
        from docling.datamodel.base_models import ConversionStatus
        from docling.datamodel.document import DocumentStream

        try:
            stream = DocumentStream(name=source_name, stream=io.BytesIO(raw_content))
            result = self.converter.convert(stream)
        except Exception as exc:
            raise DocumentParseError(
                details={"source_name": source_name, "error": str(exc)}
            ) from exc

        if result.status == ConversionStatus.FAILURE:
            raise DocumentParseError(
                details={"source_name": source_name, "status": result.status.value}
            )

        return result


class LightweightParser:
    """Plain-text extraction using pypdf, python-docx, python-pptx, openpyxl.

    Used when docling is not installed. No ML layout analysis or table detection.
    """

    def to_markdown(self, source_name: str, raw_content: bytes) -> str:
        suffix = Path(source_name).suffix.lower()
        if suffix == ".pdf":
            return self._pdf(source_name, raw_content)
        if suffix == ".docx":
            return self._docx(source_name, raw_content)
        if suffix == ".pptx":
            return self._pptx(source_name, raw_content)
        raise DocumentParseError(details={"source_name": source_name})

    def to_csv_tables(self, source_name: str, raw_content: bytes) -> str:
        return self._xlsx(source_name, raw_content)

    def _pdf(self, source_name: str, raw_content: bytes) -> str:
        from pypdf import PdfReader

        try:
            reader = PdfReader(io.BytesIO(raw_content))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception as exc:
            raise DocumentParseError(
                details={"source_name": source_name, "error": str(exc)}
            ) from exc

    def _docx(self, source_name: str, raw_content: bytes) -> str:
        from docx import Document

        try:
            doc = Document(io.BytesIO(raw_content))
            return "\n".join(p.text for p in doc.paragraphs if p.text)
        except Exception as exc:
            raise DocumentParseError(
                details={"source_name": source_name, "error": str(exc)}
            ) from exc

    def _pptx(self, source_name: str, raw_content: bytes) -> str:
        from pptx import Presentation

        try:
            prs = Presentation(io.BytesIO(raw_content))
            texts = [
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            ]
            return "\n".join(texts)
        except Exception as exc:
            raise DocumentParseError(
                details={"source_name": source_name, "error": str(exc)}
            ) from exc

    def _xlsx(self, source_name: str, raw_content: bytes) -> str:
        from openpyxl import load_workbook

        try:
            wb = load_workbook(io.BytesIO(raw_content), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    rows.append(",".join("" if c is None else str(c) for c in row))
            return "\n".join(rows)
        except Exception as exc:
            raise DocumentParseError(
                details={"source_name": source_name, "error": str(exc)}
            ) from exc


binary_parser: DoclingParser | LightweightParser = (
    DoclingParser() if _DOCLING_AVAILABLE else LightweightParser()
)
